from __future__ import annotations

import argparse
import json
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable


@dataclass(frozen=True)
class AzureOpenAIConfig:
    endpoint: str
    api_key: str
    deployment: str
    api_version: str


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def _list_asset_filenames(assets_dir: Path) -> list[str]:
    if not assets_dir.exists():
        return []
    return sorted([p.name for p in assets_dir.iterdir() if p.is_file()])


def _get_azure_openai_config(
    endpoint: str | None,
    api_key: str | None,
    deployment: str | None,
    api_version: str | None,
) -> AzureOpenAIConfig:
    # Load env vars from a local .env next to this script (optional).
    # This keeps secrets out of git while avoiding manual shell exports.
    try:
        from dotenv import load_dotenv

        script_dir = Path(__file__).resolve().parent
        # Preferred location: Chapter03/Activities/.env
        load_dotenv(dotenv_path=script_dir / ".env", override=False)
        # Also support the common “assets/.env” mistake for convenience.
        load_dotenv(dotenv_path=script_dir / "assets" / ".env", override=False)
    except Exception:
        # If python-dotenv isn't installed, continue with normal env var resolution.
        pass

    resolved_endpoint = endpoint or os.getenv("AZURE_OPENAI_ENDPOINT")
    resolved_api_key = api_key or os.getenv("AZURE_OPENAI_API_KEY")
    resolved_deployment = deployment or os.getenv("AZURE_OPENAI_DEPLOYMENT")
    resolved_api_version = api_version or os.getenv("AZURE_OPENAI_API_VERSION")

    missing: list[str] = []
    if not resolved_endpoint:
        missing.append("AZURE_OPENAI_ENDPOINT")
    if not resolved_api_key:
        missing.append("AZURE_OPENAI_API_KEY")
    if not resolved_deployment:
        missing.append("AZURE_OPENAI_DEPLOYMENT")
    if not resolved_api_version:
        missing.append("AZURE_OPENAI_API_VERSION")

    if missing:
        raise SystemExit(
            "Missing Azure OpenAI configuration. Provide via flags or env vars: "
            + ", ".join(missing)
        )

    return AzureOpenAIConfig(
        endpoint=resolved_endpoint,
        api_key=resolved_api_key,
        deployment=resolved_deployment,
        api_version=resolved_api_version,
    )


def _extract_json_object(text: str) -> dict[str, Any]:
    """Best-effort extraction of a JSON object from model output."""
    try:
        obj = json.loads(text)
        if isinstance(obj, dict):
            return obj
    except json.JSONDecodeError:
        pass

    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise ValueError("No JSON object found in model output")

    candidate = text[start : end + 1]
    obj = json.loads(candidate)
    if not isinstance(obj, dict):
        raise ValueError("Extracted JSON was not an object")
    return obj


def _sanitize_filename(name: str) -> str:
    name = name.strip().lower()
    name = re.sub(r"[^a-z0-9\-_.]+", "-", name)
    name = re.sub(r"-+", "-", name)
    return name.strip("-") or "slide"


def _iter_str_list(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [str(v) for v in value]
    return [str(value)]


def generate_slides_json(
    *,
    markdown: str,
    asset_filenames: list[str],
    config: AzureOpenAIConfig,
    temperature: float,
    max_slides: int,
) -> dict[str, Any]:
    from openai import AzureOpenAI

    system = (
        "You create concise training slide content. "
        "Return ONLY valid JSON (no markdown, no commentary). "
        "Keep bullets short (<= 10 words each) and slide-friendly."
    )

    user = {
        "task": "Generate a slide deck outline for this activity.",
        "constraints": {
            "max_slides": max_slides,
            "max_bullets_per_slide": 6,
            "max_words_per_bullet": 10,
            "tone": "professional, facilitated workshop",
        },
        "expected_json_schema": {
            "deck_title": "string",
            "slides": [
                {
                    "id": "string",
                    "title": "string",
                    "subtitle": "string (optional)",
                    "duration_minutes": "number (optional)",
                    "bullets": ["string"],
                    "speaker_notes": "string (optional)",
                    "suggested_asset_filename": "string (optional, e.g. slide-01-think.png)",
                }
            ],
        },
        "available_assets_in_repo": asset_filenames,
        "activity_markdown": markdown,
    }

    client = AzureOpenAI(
        azure_endpoint=config.endpoint,
        api_key=config.api_key,
        api_version=config.api_version,
    )

    messages = [
        {"role": "system", "content": system},
        {"role": "user", "content": json.dumps(user, ensure_ascii=False)},
    ]

    def _create_chat_completion(*, include_temperature: bool, include_response_format: bool):
        kwargs: dict[str, Any] = {
            "model": config.deployment,
            "messages": messages,
        }
        if include_temperature:
            kwargs["temperature"] = temperature
        if include_response_format:
            kwargs["response_format"] = {"type": "json_object"}
        return client.chat.completions.create(**kwargs)

    # Some newer model deployments only support default temperature.
    # Strategy:
    # 1) Try with response_format + temperature
    # 2) If response_format unsupported, retry without it
    # 3) If temperature unsupported, retry without it
    try:
        resp = _create_chat_completion(include_temperature=True, include_response_format=True)
    except TypeError:
        resp = _create_chat_completion(include_temperature=True, include_response_format=False)
    except Exception as exc:
        msg = str(exc)
        if "temperature" in msg and "Only the default (1) value is supported" in msg:
            try:
                resp = _create_chat_completion(
                    include_temperature=False, include_response_format=True
                )
            except TypeError:
                resp = _create_chat_completion(
                    include_temperature=False, include_response_format=False
                )
        else:
            raise

    content = resp.choices[0].message.content or ""
    data = _extract_json_object(content)

    slides = data.get("slides")
    if not isinstance(slides, list):
        raise ValueError("Model output missing 'slides' list")

    data["slides"] = slides[:max_slides]
    data.setdefault("deck_title", "Slide Deck")

    for slide in data["slides"]:
        if isinstance(slide, dict):
            slide.setdefault("bullets", [])
            slide["bullets"] = _iter_str_list(slide.get("bullets"))
            if "suggested_asset_filename" in slide and isinstance(
                slide["suggested_asset_filename"], str
            ):
                slide["suggested_asset_filename"] = _sanitize_filename(
                    slide["suggested_asset_filename"]
                )

    return data


def write_pptx(
    *,
    slides_json: dict[str, Any],
    output_path: Path,
    assets_dir: Path | None,
) -> None:
    from pptx import Presentation

    prs = Presentation()
    prs.core_properties.title = str(slides_json.get("deck_title", "Slide Deck"))

    title_slide_layout = prs.slide_layouts[0]
    title_and_content_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[6]

    deck_title = str(slides_json.get("deck_title", "Slide Deck"))

    # Title slide
    s0 = prs.slides.add_slide(title_slide_layout)
    if s0.shapes.title:
        s0.shapes.title.text = deck_title

    slides: Iterable[Any] = slides_json.get("slides", [])
    for idx, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue

        title = str(slide.get("title", f"Slide {idx}"))
        subtitle = slide.get("subtitle")
        bullets = _iter_str_list(slide.get("bullets"))
        notes = str(slide.get("speaker_notes", ""))
        asset_name = slide.get("suggested_asset_filename")

        image_path: Path | None = None
        if assets_dir and isinstance(asset_name, str) and asset_name:
            candidate = assets_dir / asset_name
            if candidate.exists():
                image_path = candidate

        if image_path:
            s = prs.slides.add_slide(blank_layout)
            # Fit image to slide
            s.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        else:
            s = prs.slides.add_slide(title_and_content_layout)
            if s.shapes.title:
                s.shapes.title.text = title

            body = s.shapes.placeholders[1].text_frame
            body.clear()

            if subtitle:
                p0 = body.paragraphs[0]
                p0.text = str(subtitle)
                p0.level = 0
                for p in bullets:
                    para = body.add_paragraph()
                    para.text = p
                    para.level = 0
            else:
                for i, p in enumerate(bullets):
                    if i == 0:
                        body.paragraphs[0].text = p
                        body.paragraphs[0].level = 0
                    else:
                        para = body.add_paragraph()
                        para.text = p
                        para.level = 0

        if notes:
            s.notes_slide.notes_text_frame.text = notes

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Generate slide content (JSON + PPTX) for an activity using Azure OpenAI."
    )
    parser.add_argument(
        "--activity-md",
        type=Path,
        default=Path(__file__).with_name("CHAPTER03_ACTIVITY01.md"),
        help="Path to the activity markdown file.",
    )
    parser.add_argument(
        "--assets-dir",
        type=Path,
        default=Path(__file__).with_name("assets"),
        help="Assets directory (used to list available assets and optionally embed images in PPTX).",
    )
    parser.add_argument(
        "--out-json",
        type=Path,
        default=Path(__file__).with_name("slides.json"),
        help="Output JSON path.",
    )
    parser.add_argument(
        "--out-pptx",
        type=Path,
        default=Path(__file__).with_name("slides.pptx"),
        help="Output PowerPoint (.pptx) path.",
    )
    parser.add_argument("--endpoint", type=str, default=None)
    parser.add_argument("--api-key", type=str, default=None)
    parser.add_argument("--deployment", type=str, default=None)
    parser.add_argument("--api-version", type=str, default=None)
    parser.add_argument(
        "--temperature",
        type=float,
        default=1.0,
        help="Sampling temperature. Some model deployments only support the default value (1.0).",
    )
    parser.add_argument("--max-slides", type=int, default=12)

    args = parser.parse_args()

    config = _get_azure_openai_config(
        endpoint=args.endpoint,
        api_key=args.api_key,
        deployment=args.deployment,
        api_version=args.api_version,
    )

    markdown = _read_text(args.activity_md)
    assets = _list_asset_filenames(args.assets_dir)

    slides_json = generate_slides_json(
        markdown=markdown,
        asset_filenames=assets,
        config=config,
        temperature=args.temperature,
        max_slides=args.max_slides,
    )

    args.out_json.write_text(
        json.dumps(slides_json, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    write_pptx(
        slides_json=slides_json,
        output_path=args.out_pptx,
        assets_dir=args.assets_dir,
    )

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
